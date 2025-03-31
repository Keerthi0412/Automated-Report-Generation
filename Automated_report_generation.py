import os
import re
import pandas as pd
import nltk
from collections import Counter
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Ensure NLTK tokenizer is available
nltk.download("punkt", quiet=True)

def analyze_text(text, file_type):
    """Analyzes extracted text from any file."""
    words = re.findall(r'\b\w+\b', text.lower())
    sentences = nltk.sent_tokenize(text)
    return {
        "File Type": file_type,
        "Total Words": len(words),
        "Total Sentences": len(sentences),
        "Most Common Words": Counter(words).most_common(10),
    }

def analyze_pdf(file_path):
    """Extracts and analyzes text from a PDF file."""
    try:
        text = ""
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + " " if page.extract_text() else ""
        return analyze_text(text, "PDF") if text else "No extractable text found."
    except Exception as e:
        return str(e)

def analyze_docx(file_path):
    """Extracts and analyzes text from a Word (.docx) file."""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return analyze_text(text, "DOCX") if text else "No text found in the document."
    except Exception as e:
        return str(e)

def analyze_pptx(file_path):
    """Extracts and analyzes text from a PowerPoint (.pptx) file."""
    try:
        prs = Presentation(file_path)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        return analyze_text(text, "PPTX") if text else "No text found in the presentation."
    except Exception as e:
        return str(e)

def analyze_csv(file_path):
    """Analyzes a CSV file for basic statistics."""
    try:
        df = pd.read_csv(file_path)
        return {
            "File Type": "CSV",
            "Total Rows": len(df),
            "Total Columns": len(df.columns),
            "Column Summaries": {col: {
                "Sum": df[col].sum(),
                "Average": df[col].mean(),
                "Max": df[col].max(),
                "Min": df[col].min()
            } for col in df.select_dtypes(include=['number']).columns}
        }
    except Exception as e:
        return str(e)

def analyze_txt(file_path):
    """Reads and analyzes a text file."""
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            text = file.read()
        return analyze_text(text, "Text")
    except Exception as e:
        return str(e)

def process_file(file_path):
    """Determines the file type and calls the appropriate analysis function."""
    if not os.path.exists(file_path):
        return "Error: File not found."

    file_extension = file_path.split('.')[-1].lower()

    if file_extension == "txt":
        return analyze_txt(file_path)
    elif file_extension == "pdf":
        return analyze_pdf(file_path)
    elif file_extension == "docx":
        return analyze_docx(file_path)
    elif file_extension == "pptx":
        return analyze_pptx(file_path)
    elif file_extension == "csv":
        return analyze_csv(file_path)
    else:
        return "Unsupported file type. Please use TXT, PDF, DOCX, PPTX, or CSV."

# Get file path from user
file_path = input("Enter the full file path (e.g., C:\\Users\\YourName\\Documents\\file.pdf): ").strip()

# Process and display report
if file_path:
    print(f"\nProcessing: {file_path}\n")
    result = process_file(file_path)
    print(result)
else:
    print("No file path provided.")

